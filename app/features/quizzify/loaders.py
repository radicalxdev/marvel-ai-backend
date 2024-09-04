from langchain_core.documents import Document
from app.services.logger import setup_logger
import pandas as pd
import docx
from pptx import Presentation
from typing import List, Tuple
from io import BytesIO
from app.services.tool_registry import ToolFile
from pypdf import PdfReader
from api.error_utilities import LoaderError
import requests
from bs4 import BeautifulSoup
from langchain_community.document_loaders import YoutubeLoader
from youtube_transcript_api import YouTubeTranscriptApi

logger = setup_logger(__name__)

class BytesFileLoader:
    def __init__(self, byte_files: List[Tuple[BytesIO, str]], files):
        self.byte_files = byte_files
        self.files = files

    def load(self) -> List[Document]:
        documents = []
        index = 0
        for byte_file, file_type in self.byte_files:
            if file_type.lower() == 'pdf':
                docs = self.loadPDF(byte_file,index)
            elif file_type.lower() == 'csv':
                docs = self.loadCSV(byte_file,index)
            elif file_type.lower() == 'docx':
                docs = self.loadDOCX(byte_file)
            elif file_type.lower() == 'pptx':
                docs = self.loadPPTX(byte_file, index)
            elif file_type.lower() == 'txt':
                docs = self.loadTXT(byte_file, index)

            documents.extend(docs)
            index+=1

        return documents

    def loadPDF(self, byte_file, index):
        pdf_reader = PdfReader(byte_file)  
        file = self.files[index]
        startPage = file.startPage or 0
        endPage = file.endPage or len(pdf_reader.pages) - 1
        docs = []
        for i, page in enumerate(pdf_reader.pages[startPage:endPage+1]):
            page_content = page.extract_text()
            metadata = {"source": "pdf", "page_number": i + 1}

            doc = Document(page_content=page_content, metadata=metadata)
            docs.append(doc)

        return docs

    def loadCSV(self, byte_file, index):
        # Read the CSV data from the BytesIO object
        df = pd.read_csv(byte_file)
        file = self.files[index]
        startRow = file.startRow or 0
        endRow = file.endRow or len(df) - 1
        columns = file.columns
        if columns:
            selected_df = df.iloc[startRow:endRow+1,columns]
        else:
            selected_df = df.iloc[startRow:endRow+1,:]
        # Extract content from rows with chunk size 5 and create Document objects
        chunk_size = 5
        docs = []
        for start in range(0,len(selected_df),chunk_size):
            chunk = selected_df.iloc[start:start+chunk_size]
            content = chunk.to_string(index=False)
            metadata = {"source": "csv", "chunk_number": start+1}
            
            doc = Document(page_content=content, metadata=metadata)
            docs.append(doc)

        return docs

    def loadDOCX(self, byte_file):
        docs = []
        doc = docx.Document(byte_file)
        for i, para in enumerate(doc.paragraphs):
            page_content = para.text
            metadata = {"source": "docx", "paragraph_number": i + 1}
            document = Document(page_content=page_content, metadata=metadata)
            docs.append(document)

        return docs


    def loadPPTX(self, byte_file, index):
        pptx_reader = Presentation(byte_file)
        file = self.files[index]
        startSlide = file.startSlide or 0
        endSlide = file.endSlide or len(pptx_reader.slides)
        docs = []
        for i, slide in enumerate(pptx_reader.slides):
            if i>=startSlide and i<=endSlide:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slide_text = "\n".join(slide_text)
                metadata = {"source": "pptx", "slide_number": i + 1}

                doc = Document(page_content=slide_text, metadata=metadata)
                docs.append(doc)

        return docs

    
    def loadTXT(self, byte_file, index):
        file = self.files[index]
        startLine = file.startLine or 1
        endLine = file.endLine or 10**9
        content = []
        docs = []
        for current_line_number, line in enumerate(byte_file, start=1):
            if startLine <= current_line_number <= endLine:
                text = line.decode('utf-8')
                content.append(text)
            elif current_line_number > endLine:
                break
        result_text = ' '.join(s.strip() for s in content)
        metadata = {"source": "txt"}
        doc = Document(page_content=result_text, metadata=metadata)
        docs.append(doc)
        return docs



class WebPageLoader:
    def __init__(self, verbose):
        self.verbose = verbose

    def load(self, files) -> List[Document]:
        docs = []
        for file in files:
            try:
                url = file.url
                tag = file.tag
                response = requests.get(url)
                soup = BeautifulSoup(response.content,'html.parser')
                if tag:
                    content = soup.find_all(tag)
                    text = ' '.join(data.get_text().strip() for data in content)
                else:
                    text = soup.get_text(separator="\n", strip=True)

                metadata = {"source": "webpage"}
                doc = Document(page_content=text, metadata=metadata)
                docs.append(doc)

            except Exception as e:
                logger.error(f"An error occurred while processing web page at {file}: {str(e)}")
                raise LoaderError(f"Error loading web page: {file}") from e

        return docs


class CustomYoutubeLoader:

    def __init__(self, verbose=None):
        self.verbose = verbose

    def transcript_check(self, video_id):
        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id=video_id,
                                                             languages=['en', 'en-GB', 'en-US', 'en-AU', 'en-BZ',
                                                                        'en-CA', 'en-IE', 'en-JM', 'en-NZ', 'en-ZA',
                                                                        'en-TT', 'en-GB'])
        except:
            return False

        return transcript


    def load(self, files):

        documents = []
        for file in files:
            url = file.url
            start_timestamp = file.start_timestamp
            end_timestamp = file.end_timestamp 

            video_id = YoutubeLoader.extract_video_id(url)

            data = ''

            transcript_data = self.transcript_check(video_id)
            if transcript_data:
                try:
                    
                    # Filter transcript based on start_time and end_time
                    if start_timestamp is not None and end_timestamp is not None:
                        filtered_transcript = [entry for entry in transcript_data if start_timestamp <= entry['start'] <= end_timestamp]
                    elif start_timestamp is not None:
                        filtered_transcript = [entry for entry in transcript_data if entry['start'] >= start_timestamp]
                    elif end_timestamp is not None:
                        filtered_transcript = [entry for entry in transcript_data if entry['start'] <= end_timestamp]
                    else:
                        filtered_transcript = transcript_data
                    
                    # Extract and join the text from the filtered transcript data
                    transcript_text = ' '.join([entry['text'] for entry in filtered_transcript])

                    metadata = {'source': video_id}

                    document_data = Document(page_content=transcript_text, metadata=metadata)
                    documents.append(document_data)

                except Exception as e:
                    logger.error(f"An error occurred while processing youtube data: {str(e)}")
                    raise LoaderError(f"Error loading youtube data ") from e


        return documents


    




