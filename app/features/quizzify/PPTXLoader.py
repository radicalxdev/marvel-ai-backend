from typing import List
from io import BytesIO
from pptx import Presentation
from langchain_core.documents import Document

class PPTXLoader:
    def __init__(self):
        """
        Initializes the PPTXLoader.
        """
        pass

    def load(self, file_paths: List[str]) -> List[Document]:
        """
        Loads PPTX files and converts them into Document objects.

        :param file_paths: List of file paths representing the files to load.
        :return: List of Document objects.
        """
        documents = []
        for file_path in file_paths:
            file_type = file_path.split(".")[-1]

            if file_type.lower() not in ["ppt", "pptx"]:
                raise ValueError(f"Unsupported file type: {file_type}")

            with open(file_path, 'rb') as f:
                bytes_io_obj = BytesIO(f.read())

            ppt_content = self.read_pptx(bytes_io_obj)

            metadata = {"source": "PPTX", "file_name": file_path}
            doc = Document(page_content=ppt_content, metadata=metadata)
            documents.append(doc)

        return documents

    def read_pptx(self, bytes_io_obj: BytesIO) -> str:
        """
        Reads the content of a PPTX file from a BytesIO object.

        :param bytes_io_obj: BytesIO object containing PPTX file content.
        :return: Text content of the PPTX file.
        """
        prs = Presentation(bytes_io_obj)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

# Example usage
if __name__ == "__main__":
    file_paths = ["C:/Users/Fahira/Downloads/TLC AI in the Classroom.pptx"] 
    loader = PPTXLoader()
    documents = loader.load(file_paths)
    for doc in documents:
        print(doc.page_content)
