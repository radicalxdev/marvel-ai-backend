from typing import List, Tuple, Dict, Any
from io import BytesIO
from fastapi import UploadFile
from pypdf import PdfReader
from urllib.parse import urlparse
import requests
import os
import json
import time
from docx import Document as DocxDocument # new
from pptx import Presentation as PptxPresentation
from youtube_transcript_api import (NoTranscriptFound, TranscriptsDisabled, YouTubeTranscriptApi)
import csv # new
from bs4 import BeautifulSoup
import requests

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough, RunnableParallel
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.pydantic_v1 import BaseModel, Field
from langchain_google_genai import GoogleGenerativeAI
from langchain_google_genai import GoogleGenerativeAIEmbeddings

from app.services.logger import setup_logger
from app.services.tool_registry import ToolFile
from app.api.error_utilities import LoaderError

YOUTUBE_DOMAINS=[
    "/youtu.be/",
    "/www.youtube.com/"
]


relative_path = "features/quzzify"

logger = setup_logger(__name__)

def transform_json_dict(input_data: dict) -> dict:
    # Validate and parse the input data to ensure it matches the QuizQuestion schema
    quiz_question = QuizQuestion(**input_data)

    # Transform the choices list into a dictionary
    transformed_choices = {choice.key: choice.value for choice in quiz_question.choices}

    # Create the transformed structure
    transformed_data = {
        "question": quiz_question.question,
        "choices": transformed_choices,
        "answer": quiz_question.answer,
        "explanation": quiz_question.explanation
    }

    return transformed_data

def read_text_file(file_path):
    # Get the directory containing the script file
    script_dir = os.path.dirname(os.path.abspath(__file__))

    # Combine the script directory with the relative file path
    absolute_file_path = os.path.join(script_dir, file_path)
    
    with open(absolute_file_path, 'r') as file:
        return file.read()

class RAGRunnable:
    def __init__(self, func):
        self.func = func
    
    def __or__(self, other):
        def chained_func(*args, **kwargs):
            # Result of previous function is passed as first argument to next function
            return other(self.func(*args, **kwargs))
        return RAGRunnable(chained_func)
    
    def __call__(self, *args, **kwargs):
        return self.func(*args, **kwargs)

class UploadPDFLoader:
    def __init__(self, files: List[UploadFile]):
        self.files = files

    def load(self) -> List[Document]:
        documents = []

        for upload_file in self.files:
            with upload_file.file as pdf_file:
                pdf_reader = PdfReader(pdf_file)

                for i, page in enumerate(pdf_reader.pages):
                    page_content = page.extract_text()
                    metadata = {"source": upload_file.filename, "page_number": i + 1}

                    doc = Document(page_content=page_content, metadata=metadata)
                    documents.append(doc)

        return documents

class BytesFilePDFLoader:
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            logger.debug(file_type)
            if file_type.lower() == "pdf":
                pdf_reader = PdfReader(file) #! PyPDF2.PdfReader is deprecated

                for i, page in enumerate(pdf_reader.pages):
                    page_content = page.extract_text()
                    metadata = {"source": file_type, "page_number": i + 1}

                    doc = Document(page_content=page_content, metadata=metadata)
                    documents.append(doc)
                    
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
        return documents

class BytesFileDocxLoader:
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            if file_type.lower() in ["docx", "doc"]:
                docx_document = DocxDocument(file)
                full_text = []
                
                for paragraph in docx_document.paragraphs:
                    full_text.append(paragraph.text)
                
                page_content = "\n".join(full_text)
                metadata = {"source": file_type}
                
                doc = Document(page_content=page_content, metadata=metadata)
                documents.append(doc)
            
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        
        return documents


class PPTXLoader:
    def __init__(self, files: List[Tuple[BytesIO, str]], start: int = 2, end: int = 5):
        self.files = files
        self.start = start
        self.end = end

    def load(self) -> List[Document]:
        documents = []

        for file, file_type in self.files:
            if file_type.lower() in ["pptx", "ppt"]:
                pptx_presentation = PptxPresentation(file)
                full_text = []
                current_slide = 1

                if (self.start == None):
                    self.start = 1
                elif (self.start < 1):
                    logger.error(f"Specified starting slide is not within min range. Starting slide is now {1}")
                    self.start = 1

                if (self.end == None):
                    self.end = len(pptx_presentation.slides)
                elif (self.end > len(pptx_presentation.slides)):
                    logger.error(f"Specified final slide is not within max range. Final slide is now {len(pptx_presentation.slides)}")
                    self.end = len(pptx_presentation.slides)

                logger.info(f"Selecting slides from {self.start} to {self.end}")
                for slide in pptx_presentation.slides:
                    logger.info(f"Current slide is {current_slide}")
                    if (current_slide >= self.start and current_slide <= self.end):
                        logger.info(f"Processed slide is {current_slide}")
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    full_text.append(run.text)
                    current_slide += 1                

                page_content = "\n".join(full_text)
                metadata = {"source": file_type}

                doc = Document(page_content=page_content, metadata=metadata)
                documents.append(doc)

            else:
                raise ValueError(f"Unsupported file type: {file_type}")

        return documents

class YouTubeLoader:
    def __init__(self, url:str, start_time: float = None, end_time: float = None):
        self.url = url
        self.start_time = start_time
        self.end_time = end_time

    # The only hostnames that work are youtu.be and youtube.com
    def get_video_id(self) -> str:
        parsed_url = urlparse(self.url)
        if parsed_url.hostname == 'youtu.be':
            return parsed_url.path[1:]
        if parsed_url.hostname in ('www.youtube.com', 'youtube.com'):
            if parsed_url.path == '/watch':
                p = urlparse.parse_qs(parsed_url.query)
                return p['v'][0]
            if parsed_url.path[:7] == '/embed/':
                return parsed_url.path.split('/')[2]
            if parsed_url.path[:3] == '/v/':
                return parsed_url.path.split('/')[2]
        logger.error(f"Failed to get video id")
        raise LoaderError

    def filter_by_time_stamp(self,list_of_dicts, start=None, end=None):
    # Define the filtering function based on the provided min and/or max values
        def is_within_range(d):
            if start is not None and float(d['start']) < float(start):
                return False
            if end is not None and float(d['start']) > float(end):
                return False
            return True
        # Apply the filtering function to the list
        logger.info(f"Filtering the data between {start} and {end}")
        return [d for d in list_of_dicts if is_within_range(d)]

    def load(self) -> List[Document]:
        full_transcript = []

        video_id = self.get_video_id()
        metadata = {"source": video_id}
        logger.info(f"Successfully gathered video id {video_id}")

        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
        except NoTranscriptFound:
            raise LoaderError(f"Failed to get transcript")

        if (self.start_time is not None or self.end_time is not None):
            # Time Stamp Retrieval
            filetred_transcript_peices =[]
            filetred_transcript_peices = self.filter_by_time_stamp(list_of_dicts=transcript,start=self.start_time,end=self.end_time)

            # check if there is any transcripts within the time stamps
            if not len(filetred_transcript_peices) > 0:
                raise ValueError(f"No video transcripts available for given time stamps")

            full_transcript = " ".join([t["text"].strip(" ") for t in filetred_transcript_peices])
        # If there is not a start or end time it uses the full transcript
        else:
            full_transcript = " ".join([t["text"].strip(" ") for t in transcript])

        return [Document(page_content=full_transcript, metadata=metadata)]

class BytesFileCSVLoader:
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            if file_type.lower() == "csv":
                csv_content = self._read_csv(file)
                metadata = {"source": file_type}
                
                doc = Document(page_content=csv_content, metadata=metadata)
                documents.append(doc)
            
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        
        return documents
    
    def _read_csv(self, file: BytesIO) -> str:
        file.seek(0)
        csv_reader = csv.reader(file.read().decode('utf-8').splitlines())
        csv_content = []
        
        for row in csv_reader:
            csv_content.append(", ".join(row))
        
        return "\n".join(csv_content)

class BytesFileTxtLoader:
    def __init__(self, files: List[Tuple[BytesIO, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for file, file_type in self.files:
            if file_type.lower() == "txt":
                txt_content = self._read_txt(file)
                metadata = {"source": file_type}
                
                doc = Document(page_content=txt_content, metadata=metadata)
                documents.append(doc)
            
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        
        return documents
    
    def _read_txt(self, file: BytesIO) -> str:
        file.seek(0)
        txt_content = file.read().decode('utf-8')
        return txt_content

class BytesFileWebPageLoader:
    def __init__(self, files: List[Tuple[str, str]]):
        self.files = files
    
    def load(self) -> List[Document]:
        documents = []
        
        for url, file_type in self.files:
            page_content = self._fetch_web_page(url)
            metadata = {"source": "web", "url": url}
            doc = Document(page_content=page_content, metadata=metadata)
            documents.append(doc)
        
        return documents
    
    def _fetch_web_page(self, url: str) -> str:
        response = requests.get(url)
        response.raise_for_status()  # Check for request errors

        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Extract the main content. Customize this based on the structure of the target website.
        content = soup.find('main')  # 'main' tag is common for main content. Adjust as needed.
        
        if not content:
            # Fallback if 'main' tag is not found.
            content = soup.find('body')
        
        text_content = content.get_text(separator="\n") if content else "Content not found"
        
        return text_content
    
class LocalFileLoader:
    def __init__(self, file_paths: list[str], expected_file_type="pdf"):
        self.file_paths = file_paths
        self.expected_file_type = expected_file_type

    def load(self) -> List[Document]:
        documents = []
        
        # Ensure file paths is a list
        self.file_paths = [self.file_paths] if isinstance(self.file_paths, str) else self.file_paths
    
        for file_path in self.file_paths:
            
            file_type = file_path.split(".")[-1]

            if file_type != self.expected_file_type:
                raise ValueError(f"Expected file type: {self.expected_file_type}, but got: {file_type}")

            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)

                for i, page in enumerate(pdf_reader.pages):
                    page_content = page.extract_text()
                    metadata = {"source": file_path, "page_number": i + 1}

                    doc = Document(page_content=page_content, metadata=metadata)
                    documents.append(doc)

        return documents

class URLLoader:
    def __init__(self, file_loader=None, expected_file_types=["pdf","docx","doc","pptx","ppt","csv","txt","web"], verbose=False):
        self.Docxloader = BytesFileDocxLoader
        self.Pdfloader = BytesFilePDFLoader
        self.Pptxloader = PPTXLoader
        self.Youtubeloader = YouTubeLoader
        self.Csvloader = BytesFileCSVLoader
        self.Txtloader = BytesFileTxtLoader
        self.Webloader = BytesFileWebPageLoader  # Assuming this loader expects a list of URLs
        self.expected_file_types = expected_file_types
        self.verbose = verbose
        self.loader = None
    
    def load(self, tool_files: List[ToolFile]) -> List[Document]:
        queued_files = []
        queued_urls = []  # For URLs that are web pages
        documents = []
        any_success = False

        expected_file_extensions = tuple(self.expected_file_types)

        for tool_file in tool_files:
            try:
                url = tool_file.url
                response = requests.get(url)
                parsed_url = urlparse(url)
                path = parsed_url.path

                if path.endswith(expected_file_extensions):
                    if response.status_code == 200:
                        # Read file
                        file_content = BytesIO(response.content)

                        # Check file type
                        file_type = path.split(".")[-1].lower()

                        if file_type in self.expected_file_type:
                            logger.info(f"Successfully loaded file from: {file_type}")
                            # Append to Queue
                            queued_files.append((file_content, file_type))
                        elif not url.__contains__(YOUTUBE_DOMAINS[0]) or not url.__contains__(YOUTUBE_DOMAINS[1]):
                            logger.info(f"Successfully loaded web page: {url}")
                            # Append to URL queue for web page loader
                            queued_urls.append((url, "web"))

                        if self.verbose:
                            logger.info(f"Successfully loaded file from {url}")

                        any_success = True  # Mark that at least one file was successfully loaded
                    else:
                        logger.error(f"Request failed to load file from {url} and got status code {response.status_code}")

                elif url.__contains__(YOUTUBE_DOMAINS[0]) or url.__contains__(YOUTUBE_DOMAINS[1]):
                    youtube_loader = YouTubeLoader(url=url)
                    youtube_documents = youtube_loader.load()
                    logger.info(f"Loaded youtube document")
                    documents.extend(youtube_documents)
                    logger.info(f"Added to documents")
                    any_success = True

            except Exception as e:
                logger.error(f"Failed to load file from {url}")
                logger.error(e)
                continue

        # Pass Queue to the file loader if there are any successful loads
        if any_success and queued_files:
            # Process the queued files using the appropriate loaders
            for file_content, file_type in queued_files:
                if file_type == "pdf":
                    loader = self.Pdfloader([(file_content, file_type)])
                elif file_type == "docx":
                    loader = self.Docxloader([(file_content, file_type)])
                elif file_type == "docx":
                    loader = self.Pptxloader([(file_content, file_type)])
                elif file_type == "csv":
                    loader = self.Csvloader([(file_content, file_type)])
                elif file_type == "txt":
                    loader = self.Txtloader([(file_content, file_type)])
                else:
                    continue  # Unsupported file type
                
                documents.extend(loader.load())

        if queued_urls:
            # Process the web page URLs using the web page loader
            web_loader = self.Webloader(files=queued_urls)
            documents.extend(web_loader.load())

        if self.verbose:
            logger.info(f"Loaded {len(documents)} documents")

        if not any_success:
            raise LoaderError("Unable to load any files from URLs")

        return documents

class RAGpipeline:
    def __init__(self, loader=None, splitter=None, vectorstore_class=None, embedding_model=None, verbose=False):
        default_config = {
            "loader": URLLoader(verbose = verbose), # Creates instance on call with verbosity
            "splitter": RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100),
            "vectorstore_class": Chroma,
            "embedding_model": GoogleGenerativeAIEmbeddings(model='models/embedding-001')
        }
        self.loader = loader or default_config["loader"]
        self.splitter = splitter or default_config["splitter"]
        self.vectorstore_class = vectorstore_class or default_config["vectorstore_class"]
        self.embedding_model = embedding_model or default_config["embedding_model"]
        self.verbose = verbose

    def load_PDFs(self, files) -> List[Document]:
        if self.verbose:
            logger.info(f"Loading {len(files)} files")
            logger.info(f"Loader type used: {type(self.loader)}")
        
        logger.debug(f"Loader is a: {type(self.loader)}")
        
        try:
            total_loaded_files = self.loader.load(files)
        except LoaderError as e:
            logger.error(f"Loader experienced error: {e}")
            raise LoaderError(e)
            
        return total_loaded_files
    
    def split_loaded_documents(self, loaded_documents: List[Document]) -> List[Document]:
        if self.verbose:
            logger.info(f"Splitting {len(loaded_documents)} documents")
            logger.info(f"Splitter type used: {type(self.splitter)}")
            
        total_chunks = []
        chunks = self.splitter.split_documents(loaded_documents)
        total_chunks.extend(chunks)
        
        if self.verbose: logger.info(f"Split {len(loaded_documents)} documents into {len(total_chunks)} chunks")
        
        return total_chunks
    
    def create_vectorstore(self, documents: List[Document]):
        if self.verbose:
            logger.info(f"Creating vectorstore from {len(documents)} documents")
        
        self.vectorstore = self.vectorstore_class.from_documents(documents, self.embedding_model) # here is the issue

        if self.verbose: logger.info(f"Vectorstore created")
        return self.vectorstore
    
    def compile(self):
        # Compile the pipeline
        self.load_PDFs = RAGRunnable(self.load_PDFs)
        self.split_loaded_documents = RAGRunnable(self.split_loaded_documents)
        self.create_vectorstore = RAGRunnable(self.create_vectorstore)
        if self.verbose: logger.info(f"Completed pipeline compilation")
    
    def __call__(self, documents):
        # Returns a vectorstore ready for usage 
        
        if self.verbose: 
            logger.info(f"Executing pipeline")
            logger.info(f"Start of Pipeline received: {len(documents)} documents of type {type(documents[0])}")
        
        pipeline = self.load_PDFs | self.split_loaded_documents | self.create_vectorstore
        return pipeline(documents)

class QuizBuilder:
    def __init__(self, vectorstore, topic, prompt=None, model=None, parser=None, verbose=False):
        default_config = {
            "model": GoogleGenerativeAI(model="gemini-1.0-pro"),
            "parser": JsonOutputParser(pydantic_object=QuizQuestion),
            "prompt": read_text_file("prompt/quizzify-prompt.txt")
        }
        
        self.prompt = prompt or default_config["prompt"]
        self.model = model or default_config["model"]
        self.parser = parser or default_config["parser"]
        
        self.vectorstore = vectorstore
        self.topic = topic
        self.verbose = verbose
        
        if vectorstore is None: raise ValueError("Vectorstore must be provided")
        if topic is None: raise ValueError("Topic must be provided")
    
    def compile(self):
        # Return the chain
        prompt = PromptTemplate(
            template=self.prompt,
            input_variables=["topic"],
            partial_variables={"format_instructions": self.parser.get_format_instructions()}
        )
        
        retriever = self.vectorstore.as_retriever()
        
        runner = RunnableParallel(
            {"context": retriever, "topic": RunnablePassthrough()}
        )
        
        chain = runner | prompt | self.model | self.parser
        
        if self.verbose: logger.info(f"Chain compilation complete")
        
        return chain

    def validate_response(self, response: Dict) -> bool:
        try:
            # Assuming the response is already a dictionary
            if isinstance(response, dict):
                if 'question' in response and 'choices' in response and 'answer' in response and 'explanation' in response:
                    choices = response['choices']
                    if isinstance(choices, dict):
                        for key, value in choices.items():
                            if not isinstance(key, str) or not isinstance(value, str):
                                return False
                        return True
            return False
        except TypeError as e:
            if self.verbose:
                logger.error(f"TypeError during response validation: {e}")
            return False

    def format_choices(self, choices: Dict[str, str]) -> List[Dict[str, str]]:
        return [{"key": k, "value": v} for k, v in choices.items()]
    
    def create_questions(self, num_questions: int = 5) -> List[Dict]:
        if self.verbose: logger.info(f"Creating {num_questions} questions")
        
        if num_questions > 10:
            return {"message": "error", "data": "Number of questions cannot exceed 10"}
        
        chain = self.compile()
        
        generated_questions = []
        attempts = 0
        max_attempts = num_questions * 5  # Allow for more attempts to generate questions

        while len(generated_questions) < num_questions and attempts < max_attempts:
            response = chain.invoke(self.topic)
            if self.verbose:
                logger.info(f"Generated response attempt {attempts + 1}: {response}")

            response = transform_json_dict(response)
            # Directly check if the response format is valid
            if self.validate_response(response):
                response["choices"] = self.format_choices(response["choices"])
                generated_questions.append(response)
                if self.verbose:
                    logger.info(f"Valid question added: {response}")
                    logger.info(f"Total generated questions: {len(generated_questions)}")
            else:
                if self.verbose:
                    logger.warning(f"Invalid response format. Attempt {attempts + 1} of {max_attempts}")
            
            # Move to the next attempt regardless of success to ensure progress
            attempts += 1

        # Log if fewer questions are generated
        if len(generated_questions) < num_questions:
            logger.warning(f"Only generated {len(generated_questions)} out of {num_questions} requested questions")
        
        if self.verbose: logger.info(f"Deleting vectorstore")
        self.vectorstore.delete_collection()
        
        # Return the list of questions
        return generated_questions[:num_questions]

class QuestionChoice(BaseModel):
    key: str = Field(description="A unique identifier for the choice using letters A, B, C, or D.")
    value: str = Field(description="The text content of the choice")

class QuizQuestion(BaseModel):
    question: str = Field(description="The question text")
    choices: List[QuestionChoice] = Field(description="A list of choices for the question, each with a key and a value")
    answer: str = Field(description="The key of the correct answer from the choices list")
    explanation: str = Field(description="An explanation of why the answer is correct")

    model_config = {
        "json_schema_extra": {
            "examples": """ 
                {
                "question": "What is the capital of France?",
                "choices": [
                    {"key": "A", "value": "Berlin"},
                    {"key": "B", "value": "Madrid"},
                    {"key": "C", "value": "Paris"},
                    {"key": "D", "value": "Rome"}
                ],
                "answer": "C",
                "explanation": "Paris is the capital of France."
              }
          """
        }

      }

