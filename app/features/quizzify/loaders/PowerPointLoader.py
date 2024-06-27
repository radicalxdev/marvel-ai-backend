from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

import google.generativeai as genai
from google.generativeai import GenerativeModel
import PIL
from io import BytesIO

#Setting up the model for the AI
api_key = os.environ.get('API_KEY')

genai.configure(api_key=api_key)
multimodal_model = GenerativeModel('gemini-1.5-flash',)

filename = "DataStructures.pptx"
prs = Presentation(filename)

#Extraction of all text from slides in presentation



def get_slide_text(slides):
    text_concepts = []

    
    # Iterate over each shape in the slides collection
    for shape in slides.shapes:

        # Get the title of the slide
        title = ""
        if slides.shapes.title:
            title = slides.shapes.title.text

        texts = []
        if shape.has_text_frame:
            

            # Extract text from each paragraph in the text frame
            for paragraph in shape.text_frame.paragraphs:
                # Extract text from each run in the paragraph
                for run in paragraph.runs:
                    texts.append(run.text)
          
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            
    
            image = shape.image
            image_blob = image.blob
            
            image_file = PIL.Image.open(BytesIO(image_blob))
            
            print("Writing image in AI")
            
            response = multimodal_model.generate_content(['Describe the picture', image_file])
            
            print(response.text)
            texts.append(response.text)


        text_concepts.append(texts)

    return title, text_concepts

for slide_num, slide in enumerate(prs.slides, start = 1):
    os.mkdir(f'slide {slide_num}')
    with open(f'slide {slide_num}/slide {slide_num}.txt', 'w',encoding="utf-8") as output:

        title, text_concepts = get_slide_text(slide)

        output.write("Title: " + title + '\n\n')

        count = 0

        for texts in text_concepts:
            for(text) in texts:
                output.write(text)
                count += 1

                if(count == 20):
                    count = 0
                    output.write('\n')
