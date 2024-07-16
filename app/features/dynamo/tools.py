import csv
import json
import os
from io import BytesIO

import requests
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from fastapi import HTTPException
from pptx import Presentation
from pypdf import PdfReader
import openpyxl
import xlrd

from typing import List

from app.api.error_utilities import LoaderError, VideoTranscriptError
from app.services.logger import setup_logger
from app.services.tool_registry import ToolFile

from langchain.chains.summarize import load_summarize_chain
from langchain.prompts import PromptTemplate
from langchain_core.documents import Document
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.pydantic_v1 import BaseModel, Field
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import AzureAIDocumentIntelligenceLoader, YoutubeLoader
from langchain_google_genai import GoogleGenerativeAI

from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError

from sklearn.metrics.pairwise import cosine_similarity
from langchain_openai import OpenAIEmbeddings

AZURE_END_POINT = 'https://dynamo.cognitiveservices.azure.com/'
AZURE_API_KEY = 'azure_api_key'

gcloud_auth = 'kai-ai-backend\\app\local-auth.json'

logger = setup_logger(__name__)

# AI Model
model = GoogleGenerativeAI(model="gemini-1.0-pro")


def read_text_file(file_path):
    # Get the directory containing the script file
    script_dir = os.path.dirname(os.path.abspath(__file__))

    # Combine the script directory with the relative file path
    absolute_file_path = os.path.join(script_dir, file_path)
    
    with open(absolute_file_path, 'r') as file:
        return file.read()

class PDFSubLoader():
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        documents = []
        try:
            pdf_reader = PdfReader(self.file_content)

            for i, page in enumerate(pdf_reader.pages):
                page_content = page.extract_text()
                metadata = {"source": self.file_type, "page_number": i + 1}

                doc = Document(page_content=page_content, metadata=metadata)
                documents.append(doc)

        except Exception as e:
                logger.error(f"Failed to load file from {self.file_type} sub loader")
                logger.error(e)

        return documents

class PptxSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        documents = []
        try:
            presentation = Presentation(self.file_content)
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                page_content = "\n".join(slide_text)
                metadata = {"source": self.file_type, "slide_number": i + 1}

                doc = Document(page_content=page_content, metadata=metadata)
                documents.append(doc)

        except Exception as e:
            logger.error("Failed to load file from PPTX sub loader")
            logger.error(e)

        return documents

class TextSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        try:
            text = self.file_content.read().decode('utf-8')
            doc = Document(page_content=text, metadata={"source": self.file_type})
            return [doc]
        except Exception as e:
            logger.error(f"Failed to load file from {self.file_type} sub loader")
            logger.error(e)
            return []

class JsonSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        try:
            data = json.load(self.file_content)
            content = json.dumps(data, indent=2)
            doc = Document(page_content=content, metadata={"source": self.file_type})
            return [doc]
        except json.JSONDecodeError as e:
            logger.error(f"Failed to load JSON file from {self.file_type}: {e}")
            raise ValueError("Invalid JSON content") from e

class MarkdownSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        try:
            content = self.file_content.read().decode('utf-8')
            doc = Document(page_content=content, metadata={"source": self.file_type})
            return [doc]
        except Exception as e:
            logger.error(f"Failed to load Markdown file from {self.file_type}")
            logger.error(e)
            return []

class DocxSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        documents = []
        try:
            docx_document = DocxDocument(self.file_content)
            full_text = []
            for paragraph in docx_document.paragraphs:
                full_text.append(paragraph.text)
            content = '\n'.join(full_text)
            doc = Document(page_content=content, metadata={"source": self.file_type})
            documents.append(doc)
        except Exception as e:
            logger.error(f"Failed to load file from {self.file_type} sub loader")
            logger.error(e)
        return documents

class AzureSubLoader:
    def __init__(self, url: str, azure_endpoint: str, azure_api_key: str):
        self.url = url
        self.azure_endpoint = azure_endpoint
        self.azure_api_key = azure_api_key

    def load(self) -> List[Document]:
        documents = []
        try:
            loader = AzureAIDocumentIntelligenceLoader(
                api_endpoint=self.azure_endpoint,
                api_key=self.azure_api_key,
                url_path=self.url
            )
            documents = loader.load()
        except Exception as e:
            logger.error(f"Failed to load file from Azure sub loader")
            logger.error(e)

        return documents
    
class CsvSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        try:
            self.file_content.seek(0)
            content = self.file_content.read().decode('utf-8')
            # Validate CSV format
            csv_reader = csv.reader(content.splitlines())
            rows = list(csv_reader)
            if not rows:
                raise ValueError("Empty CSV content")
            
            doc = Document(page_content=content, metadata={"source": self.file_type})
            return [doc]
        except Exception as e:
            logger.error(f"Failed to load CSV file from {self.file_type}: {e}")
            raise ValueError("Invalid CSV content") from e

class HtmlSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        try:
            self.file_content.seek(0)
            content = self.file_content.read().decode('utf-8')
            # Validate HTML format
            soup = BeautifulSoup(content, 'html.parser')
            if not soup.body:
                raise ValueError("Invalid HTML content")

            doc = Document(page_content=soup.get_text(), metadata={"source": self.file_type})
            return [doc]
        except Exception as e:
            logger.error(f"Failed to load HTML file from {self.file_type}: {e}")
            raise ValueError("Invalid HTML content") from e
        
class GoogleSheetsSubLoader:
    def __init__(self, file_content: BytesIO, file_type: str):
        self.file_content = file_content
        self.file_type = file_type

    def load(self) -> List[Document]:
        documents = []
        try:
            if self.file_content.endswith('.xlsx'):
               
               # Load the workbook
               workbook = openpyxl.load_workbook(self.file_content)
                          
               # Iterate over all sheets
               for sheet_name in workbook.sheetnames:
                   sheet = workbook[sheet_name]
                   print(f"Sheet name: {sheet_name}")

               full_workbook =[]
               # Access cells
               for row in sheet.iter_rows(values_only=True):
                   full_workbook.append(row)

               doc = Document(page_content=full_workbook, metadata={"source": self.file_type})
               documents.append(doc)   

            elif self.file_content.endswith('.xls'):
                 
                 # Read .xls file using xlrd
                 workbook = xlrd.open_workbook(self.file_content)

                 # Iterate over all sheets
                 for sheet_idx in range(workbook.nsheets):
                     sheet = workbook.sheet_by_index(sheet_idx)
                     print(f"Sheet name: {sheet.name}")
                
                 full_data =[]
                 # Access cells
                 for row_idx in range(sheet.nrows):
                     row_data = sheet.row_values(row_idx)
                     full_data.append(row_data)
               
                 doc = Document(page_content=full_data, metadata={"source": self.file_type})
                 documents.append(doc)

        except Exception as e:
            logger.error(f"Failed to load file from {self.file_type} sub loader")
            logger.error(e)
        return documents
       

class GoogleSlidesSubLoader:
    def __init__(self, presentation_id: str, credentials_path: str, file_type: str):
        self.presentation_id = presentation_id
        self.credentials_path = credentials_path
        self.file_type = file_type

    def load(self) -> List[Document]:
        SCOPES = ['https://www.googleapis.com/auth/presentations.readonly']
        try:
            creds = service_account.Credentials.from_service_account_file(self.credentials_path, scopes=SCOPES)
            service = build('slides', 'v1', credentials=creds)

            presentation = service.presentations().get(presentationId=self.presentation_id).execute()
            slides = presentation.get('slides')

            documents = []
            for i, slide in enumerate(slides):
                slide_text = []
                for element in slide.get('pageElements', []):
                    if 'shape' in element and 'text' in element['shape']:
                        for text_run in element['shape']['text']['textElements']:
                            if 'textRun' in text_run:
                                slide_text.append(text_run['textRun']['content'])

                page_content = "\n".join(slide_text)
                metadata = {"source": self.file_type, "slide_number": i + 1}

                doc = Document(page_content=page_content, metadata=metadata)
                documents.append(doc)

            return documents

        except HttpError as e:
            logger.error(f"Failed to load Google Slides file from {self.file_type}: {e}")
            return []
        except Exception as e:
            logger.error(f"An error occurred while loading Google Slides file from {self.file_type}")
            logger.error(e)
            return []
        
class URLLoader:
    def __init__(self, file_loader=None, verbose=False):
        self.loader = file_loader
        self.verbose = verbose

    def load(self, tool_files: List[ToolFile]) -> List[Document]:
        queued_files = []
        documents_list = []
        any_success = False

        for tool_file in tool_files:
            try:
                url = tool_file.url
                file_type = tool_file.filename.split(".")[-1]

                if "docs.google.com/presentation/d/" in url:
                    file_type = "google-slides"
                    presentation_id = url.split('/d/')[1].split('/')[0]
                    queued_files.append((None, file_type, presentation_id))
                    any_success = True
                    if self.verbose:
                        logger.info(f"Successfully queued Google Slides presentation from {url}")
                else:
                    response = requests.get(url)
                    if response.status_code == 200:
                        file_content = BytesIO(response.content)
                        queued_files.append((file_content, file_type, url))
                        any_success = True
                        if self.verbose:
                            logger.info(f"Successfully loaded file from {url}")
                    else:
                        logger.error(f"Request failed to load file from {url} and got status code {response.status_code}")

            except Exception as e:
                logger.error(f"Failed to load file from {url}")
                logger.error(e)
                continue

        if any_success:
            for cur_file_content, cur_file_type, identifier in queued_files:
                if cur_file_type == "pdf":
                    self.loader = PDFSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "pptx":
                    self.loader = PptxSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "txt":
                    self.loader = TextSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "json":
                    self.loader = JsonSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "md":
                    self.loader = MarkdownSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type in ['doc', 'docx', 'xls','xlsx']:
                    self.loader = AzureSubLoader(identifier, azure_endpoint=AZURE_END_POINT, azure_api_key=AZURE_API_KEY)
                elif cur_file_type == "csv":
                    self.loader = CsvSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "html":
                    self.loader = HtmlSubLoader(cur_file_content, cur_file_type)
                elif cur_file_type == "google-slides":
                    self.loader = GoogleSlidesSubLoader(identifier, gcloud_auth, cur_file_type)
                else:
                    raise LoaderError(f"Unsupported file type: {cur_file_type}")
                document = self.loader.load()
                documents_list.extend(document)

            if self.verbose:
                logger.info(f"Loaded {len(documents_list)} documents")

        if not any_success:
            raise LoaderError("Unable to load any files from URLs")

        return documents_list

def get_youtube_doc(youtube_url: str, max_video_length=600, verbose=False) -> List[Document]:
    try:
        loader = YoutubeLoader.from_youtube_url(youtube_url, add_video_info=True)
    except Exception as e:
        logger.error(f"No such video found at {youtube_url}")
        raise VideoTranscriptError(f"No video found", youtube_url) from e
    
    try:
        docs = loader.load()
        length = docs[0].metadata["length"]
        title = docs[0].metadata["title"]
        if length > max_video_length:
            raise VideoTranscriptError(f"Video is {length} seconds long, please provide a video less than {max_video_length} seconds long", youtube_url)

        if verbose:
            logger.info(f"Found video with title: {title} and length: {length}")
            logger.info(f"Generate documents from youtube video.")
        
        return docs
    
    except Exception as e:
        logger.error(f"Video transcript might be private or unavailable in 'en' or the URL is incorrect.")
        raise VideoTranscriptError(f"No video transcripts available", youtube_url) from e

def get_document_embeddings(documents: List[Document], model) -> List:
    embeddings = []
    for doc in documents:
        embedding = model.embed_query(doc.page_content)
        embeddings.append(embedding)
    return embeddings

def select_k_nearest_neighbors(documents: List[Document], embeddings: List, K: int) -> List[Document]:
    # Compute the mean embedding to represent the overall topic
    mean_embedding = sum(embeddings) / len(embeddings)
    
    # Compute cosine similarities with the mean embedding
    similarities = cosine_similarity([mean_embedding], embeddings)[0]
    
    # Select the indices of the top K similar documents
    top_k_indices = similarities.argsort()[-K:][::-1]
    
    return [documents[i] for i in top_k_indices]

def summarize_docs(documents: List[Document], K: int = 10) -> str:
    try:
        summarize_model = GoogleGenerativeAI(model="gemini-1.5-flash")
        chain = load_summarize_chain(summarize_model, chain_type="map_reduce")
        
        # Compute embeddings for all documents
        embedding_model = OpenAIEmbeddings(model="text-embedding-ada-002")
        document_embeddings = get_document_embeddings(documents, embedding_model)
        
        # Select the K-nearest neighbors based on semantic relevance
        if len(documents) > K:
            selected_documents = select_k_nearest_neighbors(documents, document_embeddings, K)
            logger.info(f"Selected top {K} nearest neighbors for summarization.")
        else:
            selected_documents = documents
        
        # Split selected documents into batches
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
        split_docs = splitter.split_documents(selected_documents)
        
        num_batches = len(split_docs)
        logger.info(f"Number of documents: {len(selected_documents)}")
        logger.info(f"Number of batches created: {num_batches}")

        result = chain.invoke(split_docs)
    
    except Exception as e:
        logger.error(f"Summarize error: {e}")
        raise

    summary_text = result["output_text"]

    # Evaluate Summary Quality
    logger.info("Evaluating summary quality...")
    summary_length = len(summary_text.split())
    logger.info(f"Summary length (in words): {summary_length}")

    return summary_text

def load_documents(youtube_url: str, files: list[ToolFile]):
    logger.info(f'Files: {files}')
    documents = []
    if youtube_url:
        documents.extend(get_youtube_doc(youtube_url))
    if files:
        loader = URLLoader(verbose=True)
        docs = loader.load(files)
        documents.extend(docs)
    
    return documents

# Summarize chain
def summarize_transcript(youtube_url: str, max_video_length=600, verbose=False) -> str:
    try:
        loader = YoutubeLoader.from_youtube_url(youtube_url, add_video_info=True)
    except Exception as e:
        logger.error(f"No such video found at {youtube_url}")
        raise VideoTranscriptError(f"No video found", youtube_url) from e
    
    try:
        docs = loader.load()
        length = docs[0].metadata["length"]
        title = docs[0].metadata["title"]
    except Exception as e:
        logger.error(f"Video transcript might be private or unavailable in 'en' or the URL is incorrect.")
        raise VideoTranscriptError(f"No video transcripts available", youtube_url) from e
    
    splitter = RecursiveCharacterTextSplitter(
        chunk_size = 1000,
        chunk_overlap = 0
    )
    
    split_docs = splitter.split_documents(docs)
    
    full_transcript = [doc.page_content for doc in split_docs]
    full_transcript = " ".join(full_transcript)

    full_transcript = [doc.page_content for doc in split_docs]
    full_transcript = " ".join(full_transcript)

    if length > max_video_length:
        raise VideoTranscriptError(f"Video is {length} seconds long, please provide a video less than {max_video_length} seconds long", youtube_url)

    if verbose:
        logger.info(f"Found video with title: {title} and length: {length}")
        logger.info(f"Combined documents into a single string.")
        logger.info(f"Beginning to process transcript...")
    
    prompt_template = read_text_file("prompt/summarize-prompt.txt")
    summarize_prompt = PromptTemplate.from_template(prompt_template)

    summarize_model = GoogleGenerativeAI(model="gemini-1.5-flash")
    
    chain = summarize_prompt | summarize_model 
    
    return chain.invoke(full_transcript)

def generate_flashcards(summary: str, verbose=False) -> list:
    # Receive the summary from the map reduce chain and generate flashcards
    parser = JsonOutputParser(pydantic_object=Flashcard)
    
    if verbose: logger.info(f"Beginning to process flashcards from summary")
    
    template = read_text_file("prompt/dynamo-prompt.txt")
    examples = read_text_file("prompt/examples.txt")
    
    cards_prompt = PromptTemplate(
        template=template,
        input_variables=["summary", "examples"],
        partial_variables={"format_instructions": parser.get_format_instructions()}
    )
    
    cards_chain = cards_prompt | model | parser
    
    try:
        response = cards_chain.invoke({"summary": summary, "examples": examples})
    except Exception as e:
        logger.error(f"Failed to generate flashcards: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate flashcards from LLM")
    
    return response

class Flashcard(BaseModel):
    concept: str = Field(description="The concept of the flashcard") 
    definition: str = Field(description="The definition of the flashcard")