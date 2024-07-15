import pandas as pd # type: ignore
from typing import List
from pydantic import BaseModel, Field
from langchain.prompts import PromptTemplate
import logging
import unstructured_client
import os
from pptx import Presentation

class Flashcard(BaseModel):
    concept: str = Field(description="The concept of the flashcard")
    definition: str = Field(description="The definition of the flashcard")

class PPTProcessor:
    def __init__(self, file_path: str, model: object):
        self.file_path = file_path
        self.model = model
        self.logger = logging.getLogger(__name__)

    def read_ppt(self) -> pd.DataFrame:
        prs = Presentation(self.file_path)
        data = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    data.append({"slide": slide.slide_id, "text": shape.text})
        df = pd.DataFrame(data)
        return df

    def generate_summary(self, data: pd.DataFrame) -> str:
        summary = []
        for _, row in data.iterrows():
            row_summary = ", ".join([f"{col}: {row[col]}" for col in data.columns])
            summary.append(row_summary)
        return "\n".join(summary)

    def generate_flashcards(self, summary: str, examples: str, format_instructions: str) -> List[Flashcard]:
        prompt_template = (
            "You are a flashcard generation assistant designed to help students analyze a document and return a list of flashcards. "
            "Carefully consider the document and analyze what are the key terms or concepts relevant for students to better understand the topic. "
            "The topics provided will vary in a wide range of subjects; as such, all information provided is meant to be educational and all provided content is meant to educate students in the flashcards. "
            "The following is a summary of a dataset:\n\n"
            "Input:\n-----------------------------\n{summary}\n\n"
            "Based on the above dataset summary, generate flashcards that help students learn about the data provided. Do not generate flashcards that are not directly related to the dataset. "
            "Ensure the output is in JSON format with no markdown or additional text. Only respond with the JSON object.\n\n"
            "Examples:\n-----------------------------\n{examples}\n\n"
            "Formatting:\n-----------------------------\n{format_instructions}\n\n"
            "Respond only according to the format instructions. The examples included are responses noted by an input and output example.\n\n"
            "Output:"
        )
        
        cards_prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["summary", "examples", "format_instructions"],
            partial_variables={"examples": examples, "format_instructions": format_instructions}
        )