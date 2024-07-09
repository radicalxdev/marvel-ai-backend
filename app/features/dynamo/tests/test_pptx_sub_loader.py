import unittest
from io import BytesIO
from pptx import Presentation
from app.features.dynamo.tools import PptxSubLoader
from app.services.logger import setup_logger

class TestPptxSubLoader(unittest.TestCase):
    def setUp(self):
        self.logger = setup_logger(__name__)

        # Create a sample PPTX content
        self.sample_pptx_content = self.create_sample_pptx("Hello, World!")

        # Create an empty PPTX content for edge case
        self.empty_pptx_content = self.create_sample_pptx("")

    def create_sample_pptx(self, text):
        pptx_content = BytesIO()
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = text
        prs.save(pptx_content)
        pptx_content.seek(0)
        return pptx_content

    def test_load_normal_case(self):
        loader = PptxSubLoader(self.sample_pptx_content, 'pptx')
        documents = loader.load()
        self.assertTrue(len(documents) > 0)
        self.assertIn('Hello, World!', documents[0].page_content)

    def test_load_empty_pptx(self):
        loader = PptxSubLoader(self.empty_pptx_content, 'pptx')
        documents = loader.load()
        self.assertEqual(len(documents), 1)
        self.assertEqual(documents[0].page_content.strip(), '')  # Check that the content is empty

    def test_load_invalid_pptx(self):
        invalid_pptx_content = BytesIO(b"This is not a PPTX content")
        loader = PptxSubLoader(invalid_pptx_content, 'pptx')
        documents = loader.load()
        self.assertEqual(len(documents), 0)  # Expect no documents to be returned

if __name__ == '__main__':
    unittest.main()
