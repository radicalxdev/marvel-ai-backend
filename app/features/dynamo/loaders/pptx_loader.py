from pptx import Presentation
from fastapi import UploadFile
from langchain_core.documents import Document

class PPTXLoader:
    def __init__(self, files: list[UploadFile]):
        self.files = files

    def load(self) -> list:
        documents = []

        for upload_file in self.files:
            prs = Presentation(upload_file.file)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            content = "\n".join(full_text)
            metadata = {"source": upload_file.filename}
            doc = Document(page_content=content, metadata=metadata)
            documents.append(doc)

        return documents
