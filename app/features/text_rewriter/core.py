import os
from typing import Dict, Any, Union
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
from .tools import export_output, process_url_input
from google.generativeai import GenerativeModel  # Import Gemini API


def rewrite_text(
    input_data: Union[str, Dict[str, Any]],
    instructions: str,
    export_format: str = "txt",
) -> str:
    """
    Rewrite the text based on given instructions.

    Args:
        input_data (str or dict): Either raw text, uploaded file data, or URL.
        instructions (str): Instructions to rewrite the text (simplify, summarize, rephrase).
        export_format (str): Output format ('txt', 'docx', 'pdf').

    Returns:
        str: File path of exported rewritten text.
    """
    try:
        # Handle input
        if isinstance(input_data, str):
            if input_data.startswith("http"):
                text = process_url_input(input_data)  # Handle URL input
            else:
                text = input_data  # Direct text input
        elif input_data["type"] == "file":
            text = process_file_input(input_data)
        else:
            raise ValueError("Unsupported input format.")

        # Handle empty input
        if not text.strip():
            raise ValueError("Input text cannot be empty.")

        # Process rewriting based on instructions
        if instructions.lower() == "simplify":
            rewritten_text = simplify_text(text)
        elif instructions.lower() == "summarize":
            rewritten_text = summarize_text(text)
        elif instructions.lower() == "rephrase":
            rewritten_text = rephrase_text(text)
        else:
            raise ValueError("Invalid instruction provided.")

        # Export rewritten text
        output_path = export_output(rewritten_text, export_format)
        return output_path

    except Exception as e:
        raise RuntimeError(f"Error processing text: {str(e)}")


def process_file_input(file_data: Dict[str, Any]) -> str:
    """Process uploaded file data based on its type."""
    file_type = file_data["format"]
    file_path = file_data["path"]

    if file_type == "pdf":
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    elif file_type == "docx":
        doc = docx.Document(file_path)
        return " ".join([para.text for para in doc.paragraphs])
    elif file_type == "csv":
        df = pd.read_csv(file_path)
        return df.to_csv(index=False, sep=" ", header=True).strip()
    elif file_type == "ppt":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
        return text
    elif file_type == "txt":
        with open(file_path, "r") as file:
            return file.read()
    else:
        raise ValueError("Unsupported file format.")


def initialize_model():
    """Initialize Gemini API model."""
    model = GenerativeModel(model_name="gemini-pro")
    return model


def simplify_text(text: str) -> str:
    """Simplify the text using Gemini."""
    model = initialize_model()
    prompt = f"Simplify the following text: {text}"
    response = model.generate_content(prompt)
    return response.text


def summarize_text(text: str) -> str:
    """Summarize the text using Gemini."""
    model = initialize_model()
    prompt = f"Summarize the following text: {text}"
    response = model.generate_content(prompt)
    return response.text


def rephrase_text(text: str) -> str:
    """Rephrase the text using Gemini."""
    model = initialize_model()
    prompt = f"Rephrase the following text: {text}"
    response = model.generate_content(prompt)
    return response.text
