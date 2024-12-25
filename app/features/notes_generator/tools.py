import os
import csv
import re
from typing import Optional
import docx

from pdfminer.high_level import extract_text as extract_pdf
from pptx import Presentation
import xml.etree.ElementTree as ET
import pandas as pd
from youtube_transcript_api import YouTubeTranscriptApi

from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI


# Function to clean the content by removing non-ASCII characters, null bytes, and non-printable characters
def clean_content(content: str) -> str:
    content = re.sub(r'[^\x00-\x7F]+', '', content)  # Remove non-ASCII characters
    content = content.replace('\x00', '')  # Remove null bytes
    content = ''.join([char if char.isprintable() else '' for char in content])  # Remove non-printable characters
    return content

# File type handlers

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from .txt files."""
    with open(file_path, 'r') as file:
        return clean_content(file.read())

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from .pdf files."""
    return clean_content(extract_pdf(file_path))

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from .docx files."""
    doc = docx.Document(file_path)
    return clean_content('\n'.join([p.text for p in doc.paragraphs]))

def extract_text_from_csv(file_path: str) -> str:
    """Extract text from .csv files."""
    with open(file_path, 'r') as file:
        reader = csv.reader(file)
        return clean_content('\n'.join([' '.join(row) for row in reader]))

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from .pptx files."""
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return clean_content("\n".join(text))

def extract_text_from_md(file_path: str) -> str:
    """Extract text from .md (Markdown) files."""
    with open(file_path, 'r') as file:
        return clean_content(file.read())

def extract_text_from_xml(file_path: str) -> str:
    """Extract text from .xml files."""
    tree = ET.parse(file_path)
    root = tree.getroot()
    text = []
    for elem in root.iter():
        if elem.text:
            text.append(elem.text.strip())
    return clean_content("\n".join(text))

def extract_text_from_xls(file_path: str) -> str:
    """Extract text from .xls files."""
    df = pd.read_excel(file_path, sheet_name=None)
    text = []
    for sheet in df.values():
        text.append(sheet.to_string(index=False))
    return clean_content("\n".join(text))

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from .xlsx files."""
    return extract_text_from_xls(file_path)  # Same as xls for now

def extract_text_from_url(video_url: str) -> str:
    """Extract text from a YouTube video URL."""
    video_id = video_url.split("v=")[-1]
    transcript = YouTubeTranscriptApi.get_transcript(video_id)
    return clean_content(' '.join([item['text'] for item in transcript]))

# Mapping file extensions to their handler functions
file_handler_map = {
    '.txt': extract_text_from_txt,
    '.pdf': extract_text_from_pdf,
    '.docx': extract_text_from_docx,
    '.csv': extract_text_from_csv,
    '.pptx': extract_text_from_pptx,
    '.md': extract_text_from_md,
    '.xml': extract_text_from_xml,
    '.xls': extract_text_from_xls,
    '.xlsx': extract_text_from_xlsx
}

def extract_text_from_file(file_path: str) -> str:
    """Extract text from a file based on its extension using the mapping."""
    file_extension = os.path.splitext(file_path)[1].lower()

    # Check if the file extension is in our handler map
    if file_extension in file_handler_map:
        handler = file_handler_map[file_extension]()  # Instantiate the handler
        return handler(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

# Function to generate notes with RAGB from the extracted content
def generate_notes_rag(content: str, focus_topic: Optional[str], structure: str) -> str:
    # 1. Create embeddings and vector store from content (one-time or repeated)
    embeddings = OpenAIEmbeddings(openai_api_key="OPEN_AI_API_KEY")
    # This can handle large content by splitting it into documents
    docs = [content]  # If you split content into chunks, store them as multiple docs
    vectorstore = FAISS.from_texts(docs, embeddings)

    # 2. Create a QA chain using your vectorstore
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    qa_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(temperature=0.7, openai_api_key="OPEN_AI_API_KEY"),
        chain_type="stuff",
        retriever=retriever
    )

    # 3. Query the chain with your focus topic
    query = f"Summarize the content focusing on '{focus_topic}' in {structure} format"
    result = qa_chain.run(query)

    return result
